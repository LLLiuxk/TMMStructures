from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_points):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if isinstance(point, list): # Support nested bullets if needed, but simple list for now
             pass

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TMMStructures 项目进度汇报"
    subtitle.text = "力热耦合二维微结构参数化逆向生成\n汇报日期：2026年3月13日"

    # Slide 2: Project Pipeline
    add_slide(prs, "项目总体目标 (Pipeline)", [
        "1. 微结构参数化生成：定义拓扑规则，自动生成多样化微结构 (进行中)",
        "2. 力热均质化计算：快速获取有效弹性与热学性质 (✅ 已完成)",
        "3. 数据集构建与模型训练：基于扩散模型实现逆向映射 (⬜ 待开始)",
        "4. 逆向设计：给定目标性能，输出最优微结构 (⬜ 待开始)"
    ])

    # Slide 3: Module I - Homogenization
    add_slide(prs, "已完成：力热均质化 (homogenize.py)", [
        "功能：像素级有限元 (Q4单元) + 周期性边界条件 (PBC)",
        "输出：有效弹性刚度矩阵 C_eff (3x3) 与有效热导率张量 k_eff (2x2)",
        "效率：256x256 网格约 20s/张，全面向量化 (Numpy/Scipy)",
        "验证：通过纯实体、纯空洞及棋盘格等物理合理性验证"
    ])

    # Slide 4: Module II - Parametric Generation
    add_slide(prs, "已完成：微结构参数化与自动生成", [
        "架构解耦：实现“生成策略”与“底层渲染”的完全分离",
        "可变边界节点：支持四条边界动态分配节点位置与厚度",
        "拓扑优化：基于图论邻接矩阵，严格通过 BFS 剔除不连通结构",
        "多样化几何：支持直线、收缩骨干、贝塞尔曲线、半圆弧连接"
    ])

    # Slide 5: Dataset Generation Strategy
    add_slide(prs, "数据集自动生成方案", [
        "配置文件驱动：通过 dataset_config.json 灵活切换模式",
        "随机采样：快速获取大规模、高多样性训练样本",
        "网格扫描 (Grid)：高维笛卡尔积穷举所有边界分配，适合学习旋转等变特征",
        "健壮性：修复了同边连接、沙漏相交等绘图漏洞，防止组合爆炸"
    ])

    # Slide 6: Next Steps
    add_slide(prs, "下阶段计划", [
        "1. 大规模生产：启动自驱动脚本生成训练所需图像-性能配对数据",
        "2. 模型架构：搭建基于扩散模型 (Diffusion Model) 的生成式网络",
        "3. 训练验证：初步验证性质矩阵到微结构拓扑的逆向生成效果"
    ])

    output_path = "TMMStructures_Progress_Report.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    path = create_presentation()
    print(f"PPT generated successfully: {path}")
